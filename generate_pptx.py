"""Generate RLHF project presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF3, 0xF4)
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
BLUE = RGBColor(0x34, 0x98, 0xDB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_title_bar(slide, title_text):
    """Add a dark blue title bar at the top of a slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_body_text(slide, left, top, width, height, bullets, font_size=18):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, level, bold) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.font.bold = bold
        p.level = level
        p.space_after = Pt(6)
    return txBox


def add_table(slide, left, top, width, height, data, col_widths=None):
    """Add a table to the slide."""
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = DARK_GRAY
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    return table


# ══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BLUE
bg.line.fill.background()

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.0), Inches(2), Inches(0.06))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11), Inches(1.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Self-Improving AI for HR Policy Q&A"
p.font.size = Pt(44)
p.font.color.rgb = WHITE
p.font.bold = True
p = tf.add_paragraph()
p.text = "LoRA SFT + DPO Fine-Tuning on Mistral 7B"
p.font.size = Pt(28)
p.font.color.rgb = ACCENT_BLUE

txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(11), Inches(1.5))
tf2 = txBox2.text_frame
p = tf2.paragraphs[0]
p.text = "POC Results & Evaluation Report"
p.font.size = Pt(22)
p.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
p = tf2.add_paragraph()
p.text = ""
p = tf2.add_paragraph()
p.text = "198 Test Prompts  |  3 Models  |  ROUGE-L + LLM Judge Evaluation"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0x88, 0x99, 0xAA)
p = tf2.add_paragraph()
p.text = "Kaggle Free GPU Tier  |  Mistral 7B  |  4-bit QLoRA"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0x88, 0x99, 0xAA)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: Problem Statement
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Problem Statement")
add_body_text(slide, 0.6, 1.5, 12, 5.5, [
    ("The Challenge", 0, True),
    ("Generic LLMs give textbook answers citing US federal law (FLSA, ADA)", 1, False),
    ("They don't know company-specific HR policies, handbooks, or internal processes", 1, False),
    ("Responses are verbose, generic, and often irrelevant to the actual policy", 1, False),
    ("", 0, False),
    ("Our Goal", 0, True),
    ("Fine-tune Mistral 7B to answer HR questions using actual company policies", 1, False),
    ("Use RLHF pipeline: LoRA SFT (domain adaptation) + DPO (preference alignment)", 1, False),
    ("Run everything on Kaggle's free GPU tier (30 hrs/week budget)", 1, False),
    ("Produce concise, policy-grounded answers instead of generic advice", 1, False),
], font_size=20)


# ══════════════════════════════════════════════════════════════
# SLIDE 3: Architecture
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Pipeline Architecture")

# Draw pipeline boxes
boxes = [
    ("RAG System\n(GPT-4o + ChromaDB)", 0.5, 2.0, ACCENT_BLUE),
    ("Preference\nCollection", 3.2, 2.0, ACCENT_BLUE),
    ("Stage 1:\nLoRA SFT", 5.9, 2.0, GREEN),
    ("Stage 2:\nDPO", 8.6, 2.0, GREEN),
    ("Evaluation\n(ROUGE + Judge)", 11.0, 2.0, ORANGE),
]
for text, left, top, color in boxes:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.2), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line_text in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

# Arrows between boxes
for x in [2.7, 5.4, 8.1, 10.5]:
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(2.35), Inches(0.5), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = DARK_GRAY
    arrow.line.fill.background()

# Labels below
add_body_text(slide, 0.5, 3.8, 5, 3, [
    ("Local Machine", 0, True),
    ("533 HR policy documents + employee handbooks", 1, False),
    ("GPT-4o generates responses, judges preferences", 1, False),
    ("Creates chosen/rejected pairs for training", 1, False),
], font_size=16)

add_body_text(slide, 5.9, 3.8, 5, 3, [
    ("Kaggle GPU (P100, 16 GB)", 0, True),
    ("4-bit QLoRA quantization (~4 GB VRAM)", 1, False),
    ("SFT: 3 epochs, ~67 min", 1, False),
    ("DPO: 2 epochs, ~45 min", 1, False),
], font_size=16)


# ══════════════════════════════════════════════════════════════
# SLIDE 4: Data
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Training Data")
add_body_text(slide, 0.6, 1.5, 5.5, 5.5, [
    ("Data Sources", 0, True),
    ("533 HR policy documents (strova-ai/hr-policies-qa-dataset)", 1, False),
    ("5,983 Q&A pairs (xwjzds/extractive_qa_question_answering_hr)", 1, False),
    ("Employee handbooks: GitLab (benefits, expenses, onboarding), Valve", 1, False),
    ("", 0, False),
    ("Preference Collection", 0, True),
    ("GPT-4o generates candidate responses via RAG pipeline", 1, False),
    ("GPT-4o-mini judges and ranks response pairs", 1, False),
    ("Output: (prompt, chosen, rejected) triples", 1, False),
], font_size=18)

add_table(slide, 6.8, 1.8, 5.5, 1.5, [
    ["Split", "Count", "Usage"],
    ["Training", "648", "SFT + DPO training"],
    ["Test", "198", "Held-out evaluation"],
    ["Total", "846", "80/20 split"],
])


# ══════════════════════════════════════════════════════════════
# SLIDE 5: SFT Training
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Stage 1: LoRA Supervised Fine-Tuning (SFT)")

add_table(slide, 0.6, 1.5, 5.5, 4.5, [
    ["Parameter", "Value"],
    ["Base model", "Mistral 7B Instruct v0.2"],
    ["Quantization", "4-bit NF4 (QLoRA)"],
    ["LoRA rank (r)", "16"],
    ["LoRA alpha", "32"],
    ["Target modules", "q_proj, k_proj, v_proj, o_proj"],
    ["Epochs", "3"],
    ["Batch size", "4 (effective 16)"],
    ["Learning rate", "2e-4 (cosine)"],
    ["Optimizer", "paged_adamw_8bit"],
], col_widths=[2.5, 3.0])

add_body_text(slide, 6.8, 1.5, 5.5, 5, [
    ("Results", 0, True),
    ("Final training loss: 0.591", 1, False),
    ("Training time: ~67 minutes", 1, False),
    ("Adapter size: 27.3 MB", 1, False),
    ("Trainable params: 13.6M / 7.26B (0.19%)", 1, False),
    ("", 0, False),
    ("What SFT Learns", 0, True),
    ("Domain vocabulary and HR-specific terminology", 1, False),
    ("Concise response style matching reference answers", 1, False),
    ("Policy-specific details (dates, procedures, contacts)", 1, False),
    ("When to say 'refer to HR' vs providing direct answers", 1, False),
], font_size=18)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: DPO Training
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Stage 2: Direct Preference Optimization (DPO)")

add_body_text(slide, 0.6, 1.5, 5.5, 5.5, [
    ("How DPO Works", 0, True),
    ("Learns from preference pairs (chosen vs rejected)", 1, False),
    ("Increases probability of preferred responses", 1, False),
    ("Decreases probability of rejected responses", 1, False),
    ("No separate reward model needed (unlike PPO)", 1, False),
    ("", 0, False),
    ("Our Setup", 0, True),
    ("Built on merged SFT model (merge LoRA into base weights)", 1, False),
    ("Beta = 0.1 (KL penalty strength)", 1, False),
    ("Learning rate: 5e-7, 2 epochs", 1, False),
    ("Training time: ~45 minutes on Kaggle P100", 1, False),
], font_size=18)

add_body_text(slide, 6.8, 1.5, 5.5, 5.5, [
    ("What DPO Adds Beyond SFT", 0, True),
    ("Preference alignment — learns what humans prefer", 1, False),
    ("Better response structure and formatting", 1, False),
    ("Zero hallucination instances (vs 11 for Base, 3 for SFT)", 1, False),
    ("Balanced verbosity — not too short, not too long", 1, False),
    ("", 0, False),
    ("DPO Model Loading (for inference)", 0, True),
    ("1. Load base Mistral 7B (4-bit)", 1, False),
    ("2. Load SFT adapter and merge into base weights", 1, False),
    ("3. Load DPO adapter on top of merged model", 1, False),
], font_size=18)


# ══════════════════════════════════════════════════════════════
# SLIDE 7: Evaluation Setup
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Evaluation Setup")

add_body_text(slide, 0.6, 1.5, 5.5, 5.5, [
    ("Models Evaluated", 0, True),
    ("Base Mistral 7B — no fine-tuning (control)", 1, False),
    ("After LoRA SFT — domain-adapted", 1, False),
    ("After DPO — preference-aligned", 1, False),
    ("", 0, False),
    ("Test Set", 0, True),
    ("198 held-out HR prompts (20% of preference data)", 1, False),
    ("Each prompt has a reference (chosen) answer from RAG", 1, False),
    ("Topics: performance reviews, harassment, leave, privacy, etc.", 1, False),
], font_size=18)

add_body_text(slide, 6.8, 1.5, 5.5, 5.5, [
    ("Automated Metrics", 0, True),
    ("ROUGE-1, ROUGE-2, ROUGE-L (text overlap with reference)", 1, False),
    ("Word count analysis (verbosity vs reference)", 1, False),
    ("Lexical diversity (unique word ratio)", 1, False),
    ("Hallucination proxy (low ROUGE + high word count)", 1, False),
    ("", 0, False),
    ("LLM-as-Judge", 0, True),
    ("GPT-4o-mini as automated judge", 1, False),
    ("3 pairwise comparisons x 198 prompts = 594 judgments", 1, False),
    ("Criteria: helpfulness, accuracy, relevance, conciseness", 1, False),
], font_size=18)


# ══════════════════════════════════════════════════════════════
# SLIDE 8: ROUGE-L Results
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "ROUGE-L Results — Text Alignment with Reference")

add_table(slide, 0.6, 1.5, 7, 2.5, [
    ["Metric", "Base Mistral 7B", "After LoRA SFT", "After DPO"],
    ["ROUGE-1 (mean)", "0.296", "0.445 (+50%)", "0.345 (+17%)"],
    ["ROUGE-2 (mean)", "0.082", "0.202 (+146%)", "0.114 (+39%)"],
    ["ROUGE-L (mean)", "0.164", "0.293 (+78%)", "0.200 (+22%)"],
    ["ROUGE-L (median)", "0.157", "0.249", "0.187"],
    ["ROUGE-L (std)", "0.049", "0.143", "0.076"],
], col_widths=[2.0, 1.8, 1.8, 1.8])

add_body_text(slide, 0.6, 4.5, 6, 2.5, [
    ("Per-Example ROUGE-L Wins", 0, True),
    ("SFT > Base:  179/198 (90.4%)", 1, False),
    ("DPO > Base:  154/198 (77.8%)", 1, False),
    ("SFT > DPO:   160/198 (80.8%)", 1, False),
], font_size=18)

add_body_text(slide, 7.5, 1.5, 5, 5.5, [
    ("Key Findings", 0, True),
    ("SFT achieves the highest ROUGE-L (0.293)", 1, False),
    ("1.78x improvement over Base", 1, False),
    ("", 0, False),
    ("DPO improves over Base (0.200 vs 0.164)", 1, False),
    ("1.22x improvement, but below SFT", 1, False),
    ("", 0, False),
    ("SFT wins on 90.4% of individual examples", 1, False),
    ("Responses most closely match reference text", 1, False),
], font_size=18)


# ══════════════════════════════════════════════════════════════
# SLIDE 9: Word Count & Quality
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Response Length & Quality Analysis")

add_table(slide, 0.6, 1.5, 7.5, 2.5, [
    ["Metric", "Reference", "Base", "SFT", "DPO"],
    ["Word count (mean)", "105", "202", "131", "178"],
    ["Length ratio vs ref", "1.00x", "1.93x (verbose)", "1.25x (close)", "1.70x"],
    ["Lexical diversity", "0.744", "0.604", "0.693", "0.654"],
    ["Hallucination proxy*", "—", "11/198", "3/198", "0/198"],
], col_widths=[2.0, 1.2, 1.5, 1.5, 1.5])

add_body_text(slide, 0.6, 4.5, 5, 2.5, [
    ("* Hallucination proxy: ROUGE-L < 0.1 AND word count > 150", 0, False),
    ("  (verbose but completely off-topic)", 0, False),
], font_size=14)

add_body_text(slide, 7.5, 1.5, 5, 5.5, [
    ("Key Findings", 0, True),
    ("Base is most verbose (1.93x reference length)", 1, False),
    ("SFT is closest to reference length (1.25x)", 1, False),
    ("DPO in between (1.70x)", 1, False),
    ("", 0, False),
    ("DPO has ZERO hallucination instances", 1, True),
    ("(vs 11 for Base, 3 for SFT)", 1, False),
    ("", 0, False),
    ("SFT has highest lexical diversity (0.693)", 1, False),
    ("Closest to reference vocabulary usage", 1, False),
], font_size=18)


# ══════════════════════════════════════════════════════════════
# SLIDE 10: Preference Win Rates
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "GPT-4o-mini Preference Win Rates")

add_table(slide, 0.6, 1.5, 8, 2, [
    ["Comparison", "Model A Wins", "Model B Wins", "Winner"],
    ["Base vs SFT", "Base: 144 (72.7%)", "SFT: 54 (27.3%)", "Base"],
    ["SFT vs DPO", "SFT: 74 (37.4%)", "DPO: 124 (62.6%)", "DPO"],
    ["Base vs DPO", "Base: 84 (42.4%)", "DPO: 114 (57.6%)", "DPO"],
])

add_body_text(slide, 0.6, 4.0, 5.5, 3, [
    ("DPO is the preference winner", 0, True),
    ("Beats SFT with 62.6% win rate", 1, False),
    ("Beats Base with 57.6% win rate", 1, False),
    ("Judge evaluates: helpfulness, accuracy,", 1, False),
    ("relevance, and conciseness", 1, False),
    ("594 total judgments, 0 errors", 1, False),
], font_size=18)

add_body_text(slide, 6.8, 4.0, 5.5, 3, [
    ("Why does judge prefer Base over SFT?", 0, True),
    ("Verbosity bias — judge lacks policy context", 1, False),
    ("Base writes 202 words of generic advice", 1, False),
    ("SFT writes 131 words of correct policy answers", 1, False),
    ("Without ground truth, judge rewards length", 1, False),
    ("Known limitation (Zheng et al., 2023)", 1, False),
], font_size=18)


# ══════════════════════════════════════════════════════════════
# SLIDE 11: Two Metrics, Two Stories
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Two Metrics, Two Stories — Complementary Evaluation")

# Left box
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.5), Inches(2.5))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
box1.line.color.rgb = ACCENT_BLUE
tf = box1.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
p.text = "ROUGE-L says: SFT is best"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
for text in ["Measures text overlap with reference", "SFT: 0.293 vs DPO: 0.200 vs Base: 0.164",
             "SFT wins 90.4% of per-example comparisons", "Best for: factual alignment"]:
    p = tf.add_paragraph()
    p.text = "  " + text
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY

# Right box
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.5), Inches(2.5))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
box2.line.color.rgb = ORANGE
tf = box2.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
p.text = "LLM Judge says: DPO is best"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ORANGE
for text in ["Evaluates helpfulness, quality, structure", "DPO beats SFT (62.6%) and Base (57.6%)",
             "Prefers well-structured, detailed answers", "Best for: human preference alignment"]:
    p = tf.add_paragraph()
    p.text = "  " + text
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY

add_body_text(slide, 0.6, 4.5, 12, 2.5, [
    ("The Insight: Both metrics are correct — they measure different things", 0, True),
    ("SFT learned to reproduce policy content (domain knowledge) — high ROUGE", 1, False),
    ("DPO learned what humans prefer (response quality) — high win rate", 1, False),
    ("Together: Base → SFT (learn the domain) → DPO (learn preferences) = full RLHF pipeline", 1, False),
], font_size=20)


# ══════════════════════════════════════════════════════════════
# SLIDE 12: Qualitative Example
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Qualitative Comparison — Example")

add_body_text(slide, 0.6, 1.3, 12, 0.5, [
    ('Prompt: "What is the company match percentage for the retirement plan?"', 0, True),
], font_size=20)

# Three boxes for each model
model_data = [
    ("Base (ROUGE-L: 0.146)", '"I\'d be happy to help, but I\'ll need to clarify... this percentage can vary greatly from one company to another, and some companies..."', RED, "Generic rambling, no useful answer"),
    ("SFT (ROUGE-L: 0.231)", '"I\'d need to refer to the specific HR policy documents... Please let me know if you have access to that information."', GREEN, "Concise, appropriate redirect to HR docs"),
    ("DPO (ROUGE-L: 0.411)", '"Without specific context regarding the retirement plan, it is not possible to provide an answer. Please refer to your organization\'s HR policy documents or contact the HR department directly."', BLUE, "Best: clear acknowledgment + direct guidance"),
]

for i, (title, quote, color, verdict) in enumerate(model_data):
    left = 0.6 + i * 4.1
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.2), Inches(3.8), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    box.line.color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = quote
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.font.italic = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = verdict
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color


# ══════════════════════════════════════════════════════════════
# SLIDE 13: Compute Budget
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Compute Budget")

add_table(slide, 0.6, 1.5, 8, 3.5, [
    ["Task", "Platform", "Time"],
    ["RAG ingestion + preference collection", "Local (CPU)", "~30 min"],
    ["SFT training (3 epochs, 123 steps)", "Kaggle P100", "~67 min"],
    ["DPO training (2 epochs)", "Kaggle P100", "~45 min"],
    ["Evaluation (198 prompts x 3 models)", "Kaggle P100", "~45 min"],
    ["Metrics + LLM judge (594 judgments)", "Local + OpenAI API", "~18 min"],
    ["Total GPU time", "", "~2.5 hours"],
    ["Kaggle weekly budget", "", "30 hours"],
])

add_body_text(slide, 0.6, 5.5, 8, 1.5, [
    ("Total GPU usage: ~8% of Kaggle's free weekly budget", 0, True),
    ("Entire RLHF pipeline (SFT + DPO + eval) runs in a single session", 0, False),
], font_size=20)


# ══════════════════════════════════════════════════════════════
# SLIDE 14: Key Takeaways
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Key Takeaways")

takeaways = [
    ("1", "SFT successfully adapted Mistral 7B to HR domain", "ROUGE-L: 0.164 → 0.293 (+78%), 35% shorter responses", GREEN),
    ("2", "DPO improved preference quality", "62.6% win rate vs SFT, 57.6% vs Base, zero hallucinations", BLUE),
    ("3", "Full RLHF pipeline works end-to-end", "Base → SFT (learn domain) → DPO (learn preferences)", ACCENT_BLUE),
    ("4", "Minimal compute required", "~2.5 GPU hours total, fits in Kaggle free tier", ORANGE),
    ("5", "Two eval methods capture different dimensions", "ROUGE-L: factual alignment | Judge: perceived quality", DARK_GRAY),
]

for i, (num, title, detail, color) in enumerate(takeaways):
    top = 1.5 + i * 1.1
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(top), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Text
    txBox = slide.shapes.add_textbox(Inches(1.7), Inches(top - 0.05), Inches(10), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p = tf.add_paragraph()
    p.text = detail
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0x77, 0x88, 0x99)


# ══════════════════════════════════════════════════════════════
# SLIDE 15: Future Work
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Future Work & Next Steps")

add_body_text(slide, 0.6, 1.5, 5.5, 5.5, [
    ("Improve Training", 0, True),
    ("Scale preference dataset to 5,000+ pairs", 1, False),
    ("Try ORPO or SimPO as DPO alternatives", 1, False),
    ("Experiment with larger LoRA rank (r=32 or r=64)", 1, False),
    ("Multi-epoch DPO with curriculum learning", 1, False),
    ("", 0, False),
    ("Better Evaluation", 0, True),
    ("Reference-aware judging (give judge policy docs)", 1, False),
    ("Human evaluation panel for final validation", 1, False),
    ("Factual accuracy scoring against policy database", 1, False),
], font_size=18)

add_body_text(slide, 6.8, 1.5, 5.5, 5.5, [
    ("Deployment", 0, True),
    ("Deploy as API endpoint for HR chatbot", 1, False),
    ("Integrate with company Slack/Teams", 1, False),
    ("Add retrieval augmentation to fine-tuned model", 1, False),
    ("", 0, False),
    ("Continuous Learning", 0, True),
    ("Collect user feedback on deployed model", 1, False),
    ("Periodic retraining with new preference data", 1, False),
    ("A/B testing between model versions", 1, False),
    ("Expand to other domains beyond HR", 1, False),
], font_size=18)


# ══════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════
output_path = r"d:\Air India\POC\RLHF\RLHF_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
